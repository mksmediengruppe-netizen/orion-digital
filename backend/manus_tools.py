"""
ORION Digital — Manus Tools Implementation
10 новых инструментов: web_scrape, pdf_read, excel_create, slides_create,
transcribe_audio, git_execute, http_request, parallel_tasks,
research_deep, long_memory_search
"""

import os
import json
import time
import subprocess
import tempfile
import logging
from typing import Any, Dict, Optional

logger = logging.getLogger(__name__)


# ══════════════════════════════════════════════════════════════════════════
# 1. WEB SCRAPE
# ══════════════════════════════════════════════════════════════════════════

def tool_web_scrape(url: str, extract: str = "all", selector: str = None,
                    follow_links: bool = False, max_pages: int = 1) -> Dict[str, Any]:
    """Deep web scraping with BeautifulSoup."""
    try:
        import requests
        from bs4 import BeautifulSoup
        import re

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        # Disable SSL verification warnings in sandbox environments
        import urllib3
        urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

        results = {"url": url, "pages_scraped": 0, "data": {}}
        urls_to_scrape = [url]
        all_data = []

        for page_url in urls_to_scrape[:max_pages]:
            resp = requests.get(page_url, headers=headers, timeout=15, verify=False)
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "html.parser")

            # Focus on selector if provided
            root = soup.select_one(selector) if selector else soup

            page_data = {"url": page_url}

            if extract in ("all", "text"):
                # Remove scripts and styles
                for tag in root.find_all(["script", "style", "nav", "footer"]):
                    tag.decompose()
                page_data["text"] = root.get_text(separator="\n", strip=True)[:5000]

            if extract in ("all", "tables"):
                tables = []
                for table in root.find_all("table"):
                    rows = []
                    for tr in table.find_all("tr"):
                        row = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
                        if row:
                            rows.append(row)
                    if rows:
                        tables.append(rows)
                page_data["tables"] = tables

            if extract in ("all", "links"):
                links = []
                for a in root.find_all("a", href=True):
                    href = a["href"]
                    if href.startswith("http"):
                        links.append({"text": a.get_text(strip=True), "url": href})
                page_data["links"] = links[:50]

            if extract in ("all", "images"):
                images = []
                for img in root.find_all("img", src=True):
                    images.append({
                        "src": img["src"],
                        "alt": img.get("alt", "")
                    })
                page_data["images"] = images[:30]

            if extract in ("all", "contacts"):
                text = root.get_text()
                emails = list(set(re.findall(r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}", text)))
                phones = list(set(re.findall(r"[\+\(]?[1-9][0-9 .\-\(\)]{7,}[0-9]", text)))
                page_data["contacts"] = {"emails": emails[:20], "phones": phones[:20]}

            all_data.append(page_data)
            results["pages_scraped"] += 1

            # Follow pagination if requested
            if follow_links and len(urls_to_scrape) < max_pages:
                for a in soup.find_all("a", href=True):
                    href = a["href"]
                    text = a.get_text(strip=True).lower()
                    if any(kw in text for kw in ["next", "следующая", "→", "»", "2", "3"]):
                        if href.startswith("http") and href not in urls_to_scrape:
                            urls_to_scrape.append(href)

        results["data"] = all_data if len(all_data) > 1 else all_data[0] if all_data else {}
        return {"success": True, "result": results}

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 2. PDF READ
# ══════════════════════════════════════════════════════════════════════════

def tool_pdf_read(source: str, pages: str = "all", extract: str = "text",
                  question: str = None) -> Dict[str, Any]:
    """Read and analyze PDF files."""
    try:
        import pdfplumber
        import requests
        import io

        # Download if URL
        if source.startswith("http"):
            resp = requests.get(source, timeout=30)
            resp.raise_for_status()
            pdf_data = io.BytesIO(resp.content)
        else:
            pdf_data = source

        result = {}

        with pdfplumber.open(pdf_data) as pdf:
            total_pages = len(pdf.pages)
            result["total_pages"] = total_pages
            result["metadata"] = pdf.metadata if extract in ("metadata", "all") else {}

            # Parse page range
            if pages == "all":
                page_nums = list(range(total_pages))
            elif "-" in str(pages):
                start, end = pages.split("-")
                page_nums = list(range(int(start) - 1, min(int(end), total_pages)))
            elif "," in str(pages):
                page_nums = [int(p) - 1 for p in pages.split(",")]
            else:
                page_nums = [int(pages) - 1]

            all_text = []
            all_tables = []

            for i in page_nums:
                if i >= total_pages:
                    continue
                page = pdf.pages[i]

                if extract in ("text", "all"):
                    text = page.extract_text() or ""
                    all_text.append(f"--- Page {i+1} ---\n{text}")

                if extract in ("tables", "all"):
                    tables = page.extract_tables()
                    if tables:
                        all_tables.extend([{"page": i+1, "data": t} for t in tables])

            if extract in ("text", "all"):
                full_text = "\n".join(all_text)
                result["text"] = full_text[:8000]
                result["text_length"] = len(full_text)

            if extract in ("tables", "all"):
                result["tables"] = all_tables[:10]

            # Answer specific question if provided
            if question and result.get("text"):
                # Simple keyword search
                text = result["text"]
                lines = text.split("\n")
                relevant = [l for l in lines if any(
                    kw.lower() in l.lower()
                    for kw in question.split()[:5]
                )]
                result["answer_context"] = "\n".join(relevant[:20])

        return {"success": True, "result": result}

    except ImportError:
        # Fallback to pdfminer
        try:
            from pdfminer.high_level import extract_text as pdfminer_extract
            text = pdfminer_extract(source)
            return {"success": True, "result": {"text": text[:8000], "method": "pdfminer"}}
        except Exception as e2:
            return {"success": False, "error": f"pdfplumber not available, pdfminer failed: {e2}"}
    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 3. EXCEL CREATE
# ══════════════════════════════════════════════════════════════════════════

def tool_excel_create(filename: str, sheets: list, title: str = None,
                      add_charts: bool = False) -> Dict[str, Any]:
    """Create Excel files with data, formulas, and styling."""
    try:
        import openpyxl
        from openpyxl.styles import (Font, PatternFill, Alignment, Border, Side,
                                      GradientFill)
        from openpyxl.utils import get_column_letter
        from openpyxl.chart import BarChart, Reference

        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # Remove default sheet

        # Color scheme
        HEADER_FILL = PatternFill("solid", fgColor="1E3A5F")
        HEADER_FONT = Font(color="FFFFFF", bold=True, size=11)
        ALT_FILL = PatternFill("solid", fgColor="F0F4F8")
        BORDER = Border(
            left=Side(style="thin", color="CCCCCC"),
            right=Side(style="thin", color="CCCCCC"),
            top=Side(style="thin", color="CCCCCC"),
            bottom=Side(style="thin", color="CCCCCC")
        )

        output_path = f"/tmp/{filename}.xlsx"
        created_sheets = []

        for sheet_def in sheets:
            ws = wb.create_sheet(title=sheet_def.get("name", "Sheet"))
            headers = sheet_def.get("headers", [])
            data = sheet_def.get("data", [])
            formulas = sheet_def.get("formulas", {})

            # Write title if provided
            start_row = 1
            if title or sheet_def.get("name"):
                ws.cell(row=1, column=1, value=title or sheet_def.get("name"))
                ws.cell(row=1, column=1).font = Font(bold=True, size=14, color="1E3A5F")
                ws.merge_cells(start_row=1, start_column=1,
                               end_row=1, end_column=max(len(headers), 1))
                start_row = 2

            # Write headers
            if headers:
                for col_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=start_row, column=col_idx, value=header)
                    cell.fill = HEADER_FILL
                    cell.font = HEADER_FONT
                    cell.alignment = Alignment(horizontal="center", vertical="center")
                    cell.border = BORDER
                ws.row_dimensions[start_row].height = 22
                start_row += 1

            # Write data
            for row_idx, row in enumerate(data):
                fill = ALT_FILL if row_idx % 2 == 0 else PatternFill()
                for col_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=start_row + row_idx, column=col_idx, value=value)
                    cell.fill = fill
                    cell.border = BORDER
                    cell.alignment = Alignment(vertical="center")

            # Apply formulas
            for cell_ref, formula in formulas.items():
                ws[cell_ref] = formula

            # Auto-fit columns
            for col in ws.columns:
                max_len = 0
                col_letter = get_column_letter(col[0].column)
                for cell in col:
                    if cell.value:
                        max_len = max(max_len, len(str(cell.value)))
                ws.column_dimensions[col_letter].width = min(max_len + 4, 40)

            # Auto-filter
            if headers and data:
                ws.auto_filter.ref = ws.dimensions

            # Freeze header row
            ws.freeze_panes = ws.cell(row=start_row, column=1)

            # Add chart if requested
            if add_charts and data and headers:
                numeric_cols = []
                for col_idx, header in enumerate(headers, 1):
                    try:
                        float(data[0][col_idx - 1])
                        numeric_cols.append(col_idx)
                    except (ValueError, TypeError, IndexError):
                        pass

                if numeric_cols:
                    chart = BarChart()
                    chart.type = "col"
                    chart.title = sheet_def.get("name", "Chart")
                    chart.y_axis.title = "Value"
                    chart.x_axis.title = headers[0] if headers else ""
                    chart.width = 20
                    chart.height = 12

                    data_ref = Reference(ws,
                                        min_col=numeric_cols[0],
                                        max_col=numeric_cols[-1],
                                        min_row=start_row - 1,
                                        max_row=start_row + len(data) - 1)
                    chart.add_data(data_ref, titles_from_data=True)
                    ws.add_chart(chart, f"A{start_row + len(data) + 3}")

            created_sheets.append(sheet_def.get("name", "Sheet"))

        wb.save(output_path)
        file_size = os.path.getsize(output_path)

        return {
            "success": True,
            "result": {
                "path": output_path,
                "filename": f"{filename}.xlsx",
                "sheets": created_sheets,
                "size_bytes": file_size,
                "size_kb": round(file_size / 1024, 1)
            }
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 4. SLIDES CREATE
# ══════════════════════════════════════════════════════════════════════════

def tool_slides_create(filename: str, slides: list, theme: str = "modern",
                       title: str = None) -> Dict[str, Any]:
    """Create PowerPoint presentations."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Theme colors
        THEMES = {
            "modern": {
                "bg": RGBColor(0x1E, 0x3A, 0x5F),
                "accent": RGBColor(0x00, 0xB4, 0xD8),
                "text": RGBColor(0xFF, 0xFF, 0xFF),
                "content_bg": RGBColor(0xF8, 0xFA, 0xFF),
                "content_text": RGBColor(0x1E, 0x3A, 0x5F)
            },
            "dark": {
                "bg": RGBColor(0x0D, 0x0D, 0x1A),
                "accent": RGBColor(0x7C, 0x3A, 0xED),
                "text": RGBColor(0xFF, 0xFF, 0xFF),
                "content_bg": RGBColor(0x1A, 0x1A, 0x2E),
                "content_text": RGBColor(0xE0, 0xE0, 0xFF)
            },
            "clean": {
                "bg": RGBColor(0xFF, 0xFF, 0xFF),
                "accent": RGBColor(0x00, 0x78, 0xD4),
                "text": RGBColor(0x1A, 0x1A, 0x2E),
                "content_bg": RGBColor(0xF5, 0xF5, 0xF5),
                "content_text": RGBColor(0x1A, 0x1A, 0x2E)
            }
        }

        colors = THEMES.get(theme, THEMES["modern"])
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        output_path = f"/tmp/{filename}.pptx"

        def set_bg(slide, color):
            from pptx.util import Pt
            from pptx.oxml.ns import qn
            from lxml import etree
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_text_box(slide, text, left, top, width, height,
                         font_size=18, bold=False, color=None, align="left"):
            from pptx.util import Inches, Pt
            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            if color:
                p.font.color.rgb = color
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            return txBox

        for slide_def in slides:
            slide_type = slide_def.get("type", "content")
            blank_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(blank_layout)
            set_bg(slide, colors["bg"])

            slide_title = slide_def.get("title", "")
            content = slide_def.get("content", "")

            if slide_type == "title":
                # Large centered title slide
                add_text_box(slide, slide_title,
                             left=1, top=2.5, width=11.33, height=1.5,
                             font_size=44, bold=True,
                             color=colors["text"], align="center")
                subtitle = slide_def.get("subtitle", "")
                if subtitle:
                    add_text_box(slide, subtitle,
                                 left=2, top=4.2, width=9.33, height=0.8,
                                 font_size=22, color=colors["accent"], align="center")

            elif slide_type == "content":
                # Title bar at top
                add_text_box(slide, slide_title,
                             left=0.5, top=0.3, width=12.33, height=0.9,
                             font_size=28, bold=True, color=colors["text"])

                # Accent line
                from pptx.util import Inches, Pt
                line = slide.shapes.add_shape(
                    1,  # MSO_SHAPE_TYPE.RECTANGLE
                    Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.05))
                line.fill.solid()
                line.fill.fore_color.rgb = colors["accent"]
                line.line.fill.background()

                # Content area
                if content:
                    lines = content.split("\n")
                    formatted = []
                    for line_text in lines:
                        lt = line_text.strip()
                        if lt.startswith(("-", "•", "*")):
                            formatted.append(f"  • {lt.lstrip('-•* ')}")
                        elif lt:
                            formatted.append(lt)
                    content_text = "\n".join(formatted)

                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    for i, line_text in enumerate(formatted):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = line_text
                        p.font.size = Pt(18)
                        p.font.color.rgb = colors["text"]
                        p.space_after = Pt(6)

            elif slide_type == "two_col":
                add_text_box(slide, slide_title,
                             left=0.5, top=0.3, width=12.33, height=0.9,
                             font_size=28, bold=True, color=colors["text"])

                left_content = slide_def.get("left", "")
                right_content = slide_def.get("right", "")

                add_text_box(slide, left_content,
                             left=0.5, top=1.5, width=5.9, height=5.5,
                             font_size=16, color=colors["text"])
                add_text_box(slide, right_content,
                             left=6.9, top=1.5, width=5.9, height=5.5,
                             font_size=16, color=colors["text"])

            elif slide_type == "table":
                add_text_box(slide, slide_title,
                             left=0.5, top=0.3, width=12.33, height=0.9,
                             font_size=28, bold=True, color=colors["text"])

                table_data = slide_def.get("table_data", [])
                if table_data and len(table_data) > 0:
                    rows = len(table_data)
                    cols = len(table_data[0])
                    table = slide.shapes.add_table(
                        rows, cols,
                        Inches(0.5), Inches(1.5),
                        Inches(12.33), Inches(min(rows * 0.5, 5.5))
                    ).table

                    for r_idx, row in enumerate(table_data):
                        for c_idx, cell_val in enumerate(row):
                            cell = table.cell(r_idx, c_idx)
                            cell.text = str(cell_val)
                            p = cell.text_frame.paragraphs[0]
                            p.font.size = Pt(14)
                            if r_idx == 0:
                                p.font.bold = True
                                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = colors["accent"]
                            else:
                                p.font.color.rgb = colors["content_text"]

        prs.save(output_path)
        file_size = os.path.getsize(output_path)

        return {
            "success": True,
            "result": {
                "path": output_path,
                "filename": f"{filename}.pptx",
                "slides_count": len(slides),
                "theme": theme,
                "size_bytes": file_size,
                "size_kb": round(file_size / 1024, 1)
            }
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 5. TRANSCRIBE AUDIO
# ══════════════════════════════════════════════════════════════════════════

def tool_transcribe_audio(source: str, language: str = None,
                          output_format: str = "text") -> Dict[str, Any]:
    """Transcribe audio/video using Whisper."""
    try:
        import requests
        import tempfile

        # Download if URL
        local_path = source
        temp_file = None

        if source.startswith("http"):
            resp = requests.get(source, timeout=60, stream=True)
            resp.raise_for_status()
            ext = source.split(".")[-1].split("?")[0][:5]
            temp_file = tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False)
            for chunk in resp.iter_content(chunk_size=8192):
                temp_file.write(chunk)
            temp_file.close()
            local_path = temp_file.name

        # Try OpenAI Whisper API first
        api_key = os.environ.get("OPENAI_API_KEY")
        if api_key:
            try:
                import openai
                client = openai.OpenAI(api_key=api_key)
                with open(local_path, "rb") as audio_file:
                    kwargs = {
                        "model": "whisper-1",
                        "file": audio_file,
                        "response_format": "verbose_json" if output_format == "json" else "text"
                    }
                    if language:
                        kwargs["language"] = language

                    transcript = client.audio.transcriptions.create(**kwargs)

                if output_format == "json":
                    result_data = transcript.model_dump() if hasattr(transcript, 'model_dump') else str(transcript)
                else:
                    result_data = str(transcript)

                if temp_file:
                    os.unlink(temp_file.name)

                return {
                    "success": True,
                    "result": {
                        "text": result_data,
                        "language": language or "auto-detected",
                        "format": output_format,
                        "method": "openai-whisper-api"
                    }
                }
            except Exception as api_err:
                logger.warning(f"OpenAI Whisper API failed: {api_err}, trying local whisper")

        # Fallback: local whisper CLI
        cmd = ["whisper", local_path, "--output_format", output_format]
        if language:
            cmd.extend(["--language", language])
        cmd.extend(["--output_dir", "/tmp"])

        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)

        if temp_file:
            try:
                os.unlink(temp_file.name)
            except Exception:
                pass

        if result.returncode == 0:
            # Find output file
            base = os.path.splitext(os.path.basename(local_path))[0]
            out_file = f"/tmp/{base}.{output_format}"
            if os.path.exists(out_file):
                with open(out_file) as f:
                    text = f.read()
                return {
                    "success": True,
                    "result": {
                        "text": text,
                        "language": language or "auto-detected",
                        "format": output_format,
                        "method": "local-whisper"
                    }
                }
            return {"success": True, "result": {"stdout": result.stdout, "method": "local-whisper"}}
        else:
            return {"success": False, "error": f"Whisper failed: {result.stderr[:500]}"}

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 6. GIT EXECUTE
# ══════════════════════════════════════════════════════════════════════════

def tool_git_execute(operation: str, repo_url: str = None, repo_path: str = None,
                     branch: str = None, message: str = None, files: list = None,
                     pr_title: str = None, pr_body: str = None) -> Dict[str, Any]:
    """Execute Git operations."""
    try:
        def run_git(cmd, cwd=None):
            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=60, cwd=cwd
            )
            return result.returncode, result.stdout.strip(), result.stderr.strip()

        if operation == "clone":
            if not repo_url:
                return {"success": False, "error": "repo_url required for clone"}
            clone_path = repo_path or f"/tmp/{repo_url.split('/')[-1].replace('.git', '')}"
            code, out, err = run_git(["git", "clone", repo_url, clone_path])
            return {
                "success": code == 0,
                "result": {"path": clone_path, "output": out or err}
            }

        if not repo_path:
            return {"success": False, "error": "repo_path required"}

        if operation == "status":
            code, out, err = run_git(["git", "status", "--short"], repo_path)
            code2, out2, _ = run_git(["git", "log", "--oneline", "-5"], repo_path)
            return {
                "success": True,
                "result": {"status": out, "recent_commits": out2}
            }

        elif operation == "add":
            targets = files if files else ["."]
            code, out, err = run_git(["git", "add"] + targets, repo_path)
            return {"success": code == 0, "result": {"output": out or err}}

        elif operation == "commit":
            if not message:
                return {"success": False, "error": "message required for commit"}
            # Auto-add if nothing staged
            run_git(["git", "add", "-A"], repo_path)
            code, out, err = run_git(
                ["git", "commit", "-m", message], repo_path)
            return {"success": code == 0, "result": {"output": out or err}}

        elif operation == "push":
            br = branch or "main"
            code, out, err = run_git(
                ["git", "push", "origin", br], repo_path)
            return {"success": code == 0, "result": {"output": out or err}}

        elif operation == "pull":
            code, out, err = run_git(["git", "pull"], repo_path)
            return {"success": code == 0, "result": {"output": out or err}}

        elif operation == "branch":
            if branch:
                code, out, err = run_git(
                    ["git", "checkout", "-b", branch], repo_path)
            else:
                code, out, err = run_git(["git", "branch", "-a"], repo_path)
            return {"success": code == 0, "result": {"output": out or err}}

        elif operation == "log":
            code, out, err = run_git(
                ["git", "log", "--oneline", "--graph", "-20"], repo_path)
            return {"success": True, "result": {"log": out}}

        elif operation == "diff":
            code, out, err = run_git(["git", "diff", "--stat"], repo_path)
            code2, out2, _ = run_git(["git", "diff"], repo_path)
            return {
                "success": True,
                "result": {"stat": out, "diff": out2[:3000]}
            }

        elif operation == "pr":
            # Use GitHub CLI
            cmd = ["gh", "pr", "create",
                   "--title", pr_title or "New PR",
                   "--body", pr_body or ""]
            if branch:
                cmd.extend(["--head", branch])
            code, out, err = run_git(cmd, repo_path)
            return {"success": code == 0, "result": {"output": out or err}}

        return {"success": False, "error": f"Unknown operation: {operation}"}

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 7. HTTP REQUEST
# ══════════════════════════════════════════════════════════════════════════

def tool_http_request(method: str, url: str, headers: dict = None,
                      params: dict = None, body: dict = None,
                      auth_type: str = None, auth_value: str = None,
                      timeout: int = 30) -> Dict[str, Any]:
    """Make HTTP API calls."""
    try:
        import requests

        req_headers = headers or {}

        # Authentication
        auth = None
        if auth_type == "bearer" and auth_value:
            req_headers["Authorization"] = f"Bearer {auth_value}"
        elif auth_type == "api_key" and auth_value:
            req_headers["X-API-Key"] = auth_value
        elif auth_type == "basic" and auth_value:
            if ":" in auth_value:
                user, pwd = auth_value.split(":", 1)
                auth = (user, pwd)

        kwargs = {
            "headers": req_headers,
            "params": params,
            "timeout": timeout,
            "verify": True
        }
        if auth:
            kwargs["auth"] = auth
        if body:
            kwargs["json"] = body

        resp = requests.request(method, url, **kwargs)

        # Try to parse JSON response
        try:
            response_data = resp.json()
        except Exception:
            response_data = resp.text[:5000]

        return {
            "success": True,
            "result": {
                "status_code": resp.status_code,
                "ok": resp.ok,
                "headers": dict(resp.headers),
                "data": response_data,
                "url": resp.url
            }
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 8. PARALLEL TASKS
# ══════════════════════════════════════════════════════════════════════════

def tool_parallel_tasks(tasks: list, max_workers: int = 5,
                        timeout_per_task: int = 60) -> Dict[str, Any]:
    """Execute multiple Python tasks in parallel."""
    try:
        from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeout
        import traceback
        import sys
        import io

        max_workers = min(max_workers, 20)

        def execute_single(task):
            task_id = task.get("id", f"task_{tasks.index(task)}")
            code = task.get("code", "")
            start = time.time()

            # Capture stdout
            old_stdout = sys.stdout
            sys.stdout = io.StringIO()

            try:
                local_vars = {}
                exec(code, {"__builtins__": __builtins__}, local_vars)
                output = sys.stdout.getvalue()
                sys.stdout = old_stdout
                return {
                    "id": task_id,
                    "success": True,
                    "output": output,
                    "result": local_vars.get("result", None),
                    "duration_ms": round((time.time() - start) * 1000)
                }
            except Exception as e:
                sys.stdout = old_stdout
                return {
                    "id": task_id,
                    "success": False,
                    "error": str(e),
                    "traceback": traceback.format_exc()[-500:],
                    "duration_ms": round((time.time() - start) * 1000)
                }

        results = []
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(execute_single, task): task for task in tasks}
            for future, task in futures.items():
                try:
                    result = future.result(timeout=timeout_per_task)
                    results.append(result)
                except FuturesTimeout:
                    results.append({
                        "id": task.get("id", "unknown"),
                        "success": False,
                        "error": f"Timeout after {timeout_per_task}s"
                    })
                except Exception as e:
                    results.append({
                        "id": task.get("id", "unknown"),
                        "success": False,
                        "error": str(e)
                    })

        success_count = sum(1 for r in results if r.get("success"))
        return {
            "success": True,
            "result": {
                "total": len(tasks),
                "succeeded": success_count,
                "failed": len(tasks) - success_count,
                "tasks": results
            }
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 9. RESEARCH DEEP
# ══════════════════════════════════════════════════════════════════════════

def tool_research_deep(query: str, depth: int = 3, sources: int = 5,
                       output_format: str = "report",
                       progress_callback=None) -> Dict[str, Any]:
    """Deep multi-source research with AI synthesis."""
    try:
        import requests
        from bs4 import BeautifulSoup
        import openai

        depth = max(1, min(depth, 5))
        sources = max(1, min(sources, 10))

        # Step 1: Search
        if progress_callback:
            progress_callback({"step": "search", "status": "running", "detail": f"Searching: {query}"})
        search_results = []
        try:
            # Use DuckDuckGo HTML search (no API key needed)
            ddg_url = "https://html.duckduckgo.com/html/"
            resp = requests.post(ddg_url, data={"q": query}, timeout=15,
                                 headers={"User-Agent": "Mozilla/5.0"})
            soup = BeautifulSoup(resp.text, "html.parser")
            for result in soup.select(".result__body")[:sources]:
                title_el = result.select_one(".result__title")
                snippet_el = result.select_one(".result__snippet")
                link_el = result.select_one(".result__url")
                if title_el:
                    search_results.append({
                        "title": title_el.get_text(strip=True),
                        "snippet": snippet_el.get_text(strip=True) if snippet_el else "",
                        "url": link_el.get_text(strip=True) if link_el else ""
                    })
        except Exception as search_err:
            logger.warning(f"DuckDuckGo search failed: {search_err}")

        if progress_callback:
            progress_callback({"step": "search", "status": "done", "detail": f"Found {len(search_results)} sources"})

        # Step 2: Fetch content from top sources
        if progress_callback:
            progress_callback({"step": "fetch", "status": "running", "detail": f"Fetching {min(depth, len(search_results))} sources..."})
        fetched_content = []
        if depth >= 2:
            for sr in search_results[:min(depth, sources)]:
                try:
                    url = sr.get("url", "")
                    if not url.startswith("http"):
                        continue
                    resp = requests.get(url, timeout=10,
                                        headers={"User-Agent": "Mozilla/5.0"})
                    soup = BeautifulSoup(resp.text, "html.parser")
                    for tag in soup.find_all(["script", "style", "nav", "footer"]):
                        tag.decompose()
                    text = soup.get_text(separator="\n", strip=True)[:2000]
                    fetched_content.append({
                        "url": url,
                        "title": sr.get("title", ""),
                        "content": text
                    })
                except Exception:
                    pass

        if progress_callback:
            progress_callback({"step": "fetch", "status": "done", "detail": f"Fetched {len(fetched_content)} sources"})

        # Step 3: AI synthesis (via OpenRouter — same provider as main agent)
        if progress_callback:
            progress_callback({"step": "synthesis", "status": "running", "detail": "AI synthesizing findings..."})
        api_key = os.environ.get("OPENROUTER_API_KEY") or os.environ.get("OPENAI_API_KEY")
        if not api_key:
            # Return raw results without synthesis
            return {
                "success": True,
                "result": {
                    "query": query,
                    "search_results": search_results,
                    "fetched_sources": len(fetched_content),
                    "synthesis": "AI synthesis unavailable (no API key)"
                }
            }

        # Use OpenRouter if OPENROUTER_API_KEY is set, otherwise fallback to OpenAI
        base_url = "https://openrouter.ai/api/v1" if os.environ.get("OPENROUTER_API_KEY") else None
        client_kwargs = {"api_key": api_key}
        if base_url:
            client_kwargs["base_url"] = base_url
        client = openai.OpenAI(**client_kwargs)

        # Build context
        context_parts = []
        for i, sr in enumerate(search_results, 1):
            context_parts.append(f"[{i}] {sr['title']}\n{sr['snippet']}")

        for fc in fetched_content:
            context_parts.append(f"\n--- Source: {fc['title']} ---\n{fc['content']}")

        context = "\n\n".join(context_parts)[:6000]

        format_instructions = {
            "report": "Write a comprehensive research report with sections: Summary, Key Findings, Details, Conclusions.",
            "summary": "Write a concise 3-5 paragraph summary of the key findings.",
            "bullets": "List the 10 most important findings as bullet points."
        }

        prompt = f"""Research query: {query}

Sources found:
{context}

{format_instructions.get(output_format, format_instructions['report'])}

Be factual, cite sources where possible, and be thorough."""

        # Use OpenRouter model ID when going through OpenRouter
        synthesis_model = "openai/gpt-5.4-mini" if base_url else "gpt-4.1-mini"
        response = client.chat.completions.create(
            model=synthesis_model,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=2000,
            temperature=0.3
        )

        synthesis = response.choices[0].message.content

        return {
            "success": True,
            "result": {
                "query": query,
                "sources_found": len(search_results),
                "sources_analyzed": len(fetched_content),
                "synthesis": synthesis,
                "raw_sources": search_results[:5],
                "format": output_format
            }
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# 10. LONG MEMORY SEARCH
# ══════════════════════════════════════════════════════════════════════════

def tool_long_memory_search(query: str, user_id: str = None,
                            memory_type: str = "all",
                            limit: int = 10) -> Dict[str, Any]:
    """Search agent's long-term memory across sessions using memory_v9."""
    try:
        # Use memory_v9 SuperMemoryEngine
        from memory_v9 import SuperMemoryEngine
        engine = SuperMemoryEngine()
        results = engine.recall(query)
        return {
            "success": True,
            "result": {
                "query": query,
                "memories": results if isinstance(results, list) else [results] if results else [],
                "total_found": len(results) if isinstance(results, list) else (1 if results else 0),
                "memory_type": memory_type,
                "source": "memory_v9"
            }
        }
    except Exception as _mem_err:
        logger.warning(f"memory_v9 recall failed: {_mem_err}, falling back to DB search")
    # Fallback: direct DB search
    try:
        import sqlite3
        db_path = "/var/www/orion/backend/orion.db"
        if not os.path.exists(db_path):
            return {
                "success": True,
                "result": {
                    "query": query,
                    "memories": [],
                    "message": "Memory database not found"
                }
            }

        conn = sqlite3.connect(db_path)
        conn.row_factory = sqlite3.Row
        cur = conn.cursor()
        search_query = f"%{query}%"
        try:
            cur.execute(
                "SELECT * FROM memories WHERE content LIKE ? ORDER BY created_at DESC LIMIT ?",
                (search_query, limit)
            )
            rows = [dict(r) for r in cur.fetchall()]
        except Exception:
            rows = []
        conn.close()

        return {
            "success": True,
            "result": {
                "query": query,
                "memories": rows,
                "total_found": len(rows),
                "memory_type": memory_type,
                "source": "sqlite_fallback"
            }
        }

    except Exception as e:
        return {"success": False, "error": str(e)}


# ══════════════════════════════════════════════════════════════════════════
# DISPATCHER — вызывается из agent_loop.py
# ══════════════════════════════════════════════════════════════════════════

def _tool_text_to_speech_wrapper(**kwargs):
    """Wrapper for TTS tool."""
    try:
        from tts_tool import tool_text_to_speech
        return tool_text_to_speech(**kwargs)
    except ImportError:
        return {"success": False, "error": "TTS module not available"}


MANUS_TOOL_HANDLERS = {
    "web_scrape": tool_web_scrape,
    "pdf_read": tool_pdf_read,
    "excel_create": tool_excel_create,
    "slides_create": tool_slides_create,
    "transcribe_audio": tool_transcribe_audio,
    "git_execute": tool_git_execute,
    "http_request": tool_http_request,
    "parallel_tasks": tool_parallel_tasks,
    "research_deep": tool_research_deep,
    "long_memory_search": tool_long_memory_search,
    "text_to_speech": _tool_text_to_speech_wrapper,
}


def dispatch_manus_tool(tool_name: str, tool_args: dict) -> dict:
    """Dispatch a Manus tool call."""
    handler = MANUS_TOOL_HANDLERS.get(tool_name)
    if not handler:
        return {"success": False, "error": f"Unknown Manus tool: {tool_name}"}
    try:
        return handler(**tool_args)
    except TypeError as e:
        return {"success": False, "error": f"Invalid arguments for {tool_name}: {e}"}
    except Exception as e:
        logger.exception(f"Manus tool {tool_name} failed")
        return {"success": False, "error": str(e)}
